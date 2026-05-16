"""
Kuja Market Analysis Deck (v5.0) — May 2026.

A 25-slide presentation summarising the v5.0 market analysis, the
strategic pivot from a two-product play to a grant-platform-led
strategy, and the competitive landscape — with embedded charts.

Designed for partner conversations, board reviews, and internal alignment.

Regenerable: charts in docs/deck_charts.py · deck assembly here.
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

THIS_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, THIS_DIR)
import deck_charts as _charts

# Regenerate charts on every build so they stay in lockstep with content
CHART_PATHS = _charts.generate_all()
CHART = {os.path.basename(p).split("_", 1)[1].rsplit(".", 1)[0]: p
         for p in CHART_PATHS}


# ============================================================================
# BRANDING
# ============================================================================
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
SKY = RGBColor(0x2C, 0x5F, 0x8A)
CLAY = RGBColor(0xC2, 0x41, 0x0C)
CLAY_SOFT = RGBColor(0xFF, 0xE4, 0xD6)
SAND = RGBColor(0xF2, 0xF6, 0xFA)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GROW = RGBColor(0x16, 0xA3, 0x4A)
SUN = RGBColor(0xF5, 0x9E, 0x0B)
FLAG = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ============================================================================
# HELPERS
# ============================================================================

def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tx


def add_bullets(slide, left, top, width, height, items, size=14,
                color=INK, bullet_color=NAVY):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
            # Set first run for bold prefix
            if not p.runs:
                r1 = p.add_run()
            else:
                r1 = p.runs[0]
            r1.text = "• " + head
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = bullet_color
            r2 = p.add_run()
            r2.text = " " + body
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.runs[0] if p.runs else p.add_run()
            r.text = "• " + str(item)
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
        p.space_after = Pt(8)
        p.line_spacing = 1.18
    return tx


def add_rect(slide, left, top, width, height, fill, line_color=None,
             corner=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    return shape


def add_brand_bar(slide):
    """Slim brand bar at the bottom of every content slide."""
    add_rect(slide, Inches(0), Inches(7.30), SLIDE_W, Inches(0.20),
             NAVY, corner=False)
    add_text(slide, Inches(0.5), Inches(7.30), Inches(8), Inches(0.20),
             "Kuja  ·  AI-Powered Grant Management for the Global South",
             size=9, color=WHITE, align=PP_ALIGN.LEFT, italic=True)
    add_text(slide, Inches(5), Inches(7.30), Inches(7.8), Inches(0.20),
             "Market Analysis  ·  v5.0  ·  May 2026",
             size=9, color=WHITE, align=PP_ALIGN.RIGHT, italic=True)


def slide_title_bar(slide, title, subtitle=None):
    """Title banner across the top of a content slide."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85),
             NAVY, corner=False)
    add_text(slide, Inches(0.5), Inches(0.10), Inches(12.3), Inches(0.55),
             title, size=22, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.35),
                 subtitle, size=11, color=CLAY_SOFT, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_chart(slide, key, left, top, width):
    path = CHART.get(key)
    if not path or not os.path.exists(path):
        return
    slide.shapes.add_picture(path, left, top, width=width)


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============================================================================
# SLIDE 1 — Cover
# ============================================================================

s = new_slide()
# Full-bleed navy background
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY, corner=False)
# Accent stripe
add_rect(s, Inches(0), Inches(3.6), SLIDE_W, Inches(0.06), CLAY, corner=False)

# Big brand wordmark
add_text(s, Inches(0), Inches(1.7), SLIDE_W, Inches(1.6),
         "Kuja", size=110, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s, Inches(0), Inches(3.85), SLIDE_W, Inches(0.6),
         "An AI-Powered Grant Management System",
         size=24, color=CLAY_SOFT, align=PP_ALIGN.CENTER)

# Built-for line
add_text(s, Inches(0), Inches(4.55), SLIDE_W, Inches(0.45),
         "Built for the Global South",
         size=14, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

# Document type
add_text(s, Inches(0), Inches(5.7), SLIDE_W, Inches(0.5),
         "MARKET ANALYSIS",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Version
add_text(s, Inches(0), Inches(6.3), SLIDE_W, Inches(0.4),
         "v5.0  ·  May 2026",
         size=12, color=CLAY_SOFT, align=PP_ALIGN.CENTER, italic=True)

# Confidential footer
add_text(s, Inches(0), Inches(7.0), SLIDE_W, Inches(0.4),
         "INTERNAL · FOR TEAM AND PARTNERS",
         size=9, color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ============================================================================
# SLIDE 2 — Where We Started: The Original Two-Product Vision
# ============================================================================

s = new_slide()
slide_title_bar(s, "Where We Started — The Original Vision",
                "Kuja launched as a two-product play backed by a back-office service, "
                "with an aggressive sell-to-NGOs-and-donors GTM.")

# Three product columns
prod_top = Inches(1.2)
prod_h = Inches(3.6)
prod_w = Inches(4.0)
prod_gap = Inches(0.3)
prod_xs = [Inches(0.5),
           Inches(0.5) + prod_w + prod_gap,
           Inches(0.5) + 2 * (prod_w + prod_gap)]

products = [
    ("KujaLink",
     "The Marketplace",
     CLAY,
     [
         "Two-sided NGO ↔ donor discovery platform",
         "Capacity profiles, opportunity matching",
         "Originally framed as the demand-side wedge",
     ]),
    ("Kuja ERP",
     "Built on Odoo 17",
     NAVY,
     [
         "Full back-office stack: HR · Finance · Ops · Procurement",
         "Customised for NGOs with grant-aware budgeting",
         "Intended as the recurring-revenue engine",
     ]),
    ("Back-Office Services",
     "Managed Service Layer",
     SKY,
     [
         "Consultancy for NGOs without internal capacity",
         "HR, finance, operations run by Adeso team",
         "Wraparound for orgs adopting Kuja ERP",
     ]),
]

for x, (name, tag, accent, bullets) in zip(prod_xs, products):
    # Card background
    add_rect(s, x, prod_top, prod_w, prod_h, SAND)
    # Accent bar
    add_rect(s, x, prod_top, prod_w, Inches(0.55), accent, corner=False)
    add_text(s, x + Inches(0.2), prod_top + Inches(0.05),
             prod_w - Inches(0.4), Inches(0.45),
             name, size=18, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Tag line
    add_text(s, x + Inches(0.2), prod_top + Inches(0.65),
             prod_w - Inches(0.4), Inches(0.3),
             tag, size=11, italic=True, color=MUTED)
    # Bullets
    add_bullets(s, x + Inches(0.25), prod_top + Inches(1.05),
                prod_w - Inches(0.5), prod_h - Inches(1.2),
                bullets, size=11, bullet_color=accent)

# GTM band
gtm_top = Inches(5.05)
add_rect(s, Inches(0.5), gtm_top, Inches(12.3), Inches(0.8), NAVY)
add_text(s, Inches(0.7), gtm_top + Inches(0.05),
         Inches(3.0), Inches(0.7),
         "ORIGINAL GTM",
         size=12, bold=True, color=CLAY_SOFT,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(3.5), gtm_top + Inches(0.05),
         Inches(9.2), Inches(0.7),
         "Sell aggressively to NGOs AND donors  ·  "
         "monetise marketplace, ERP licences, and managed services in parallel.",
         size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Status note
note_top = Inches(6.05)
add_rect(s, Inches(0.5), note_top, Inches(12.3), Inches(1.05), CLAY_SOFT)
add_text(s, Inches(0.7), note_top + Inches(0.10),
         Inches(11.9), Inches(0.4),
         "STATUS TODAY",
         size=11, bold=True, color=CLAY)
add_text(s, Inches(0.7), note_top + Inches(0.42),
         Inches(11.9), Inches(0.55),
         "KujaLink shipped and is live in production. The ERP build is "
         "still in development after four years — over USD 3M invested. "
         "Back-office services remain a small, bespoke engagement.",
         size=12, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 3 — Investment to Date: Where the $3M+ Went
# ============================================================================

s = new_slide()
slide_title_bar(s, "Investment to Date — Where the USD 3M+ Went",
                "Four years of spend on Kuja ERP, concentrated in consultants and staff.")

# Headline number band
hl_top = Inches(1.05)
add_rect(s, Inches(0.5), hl_top, Inches(12.3), Inches(0.85), NAVY)
add_text(s, Inches(0.7), hl_top + Inches(0.05),
         Inches(6.0), Inches(0.75),
         "USD 4.23M",
         size=30, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(4.5), hl_top + Inches(0.05),
         Inches(8.0), Inches(0.75),
         "total invested in Kuja ERP, 2022 – Mar 2026",
         size=14, color=CLAY_SOFT, anchor=MSO_ANCHOR.MIDDLE)

# === LEFT: By Year ===
left_x = Inches(0.5)
left_top = Inches(2.15)
left_w = Inches(6.0)
left_h = Inches(4.7)
add_rect(s, left_x, left_top, left_w, left_h, SAND)
add_text(s, left_x + Inches(0.25), left_top + Inches(0.15),
         left_w - Inches(0.5), Inches(0.35),
         "BY YEAR (USD)", size=12, bold=True, color=NAVY)

years = [
    ("2022", 334, NAVY),
    ("2023", 1380, CLAY),
    ("2024", 1170, CLAY),
    ("2025", 1170, CLAY),
    ("2026 YTD", 166, SKY),
]
max_year = max(v for _, v, _ in years)

bar_area_left = left_x + Inches(1.4)
bar_area_top = left_top + Inches(0.65)
bar_area_w = left_w - Inches(2.0)
row_h = Inches(0.65)
row_gap = Inches(0.15)

for i, (label, value, color) in enumerate(years):
    row_top = bar_area_top + i * (row_h + row_gap)
    # Year label
    add_text(s, left_x + Inches(0.25), row_top,
             Inches(1.1), row_h,
             label, size=12, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
    # Bar
    bar_w = int(bar_area_w * (value / max_year))
    add_rect(s, bar_area_left, row_top, bar_w, row_h, color, corner=False)
    # Value label
    val_str = f"${value/1000:.2f}M" if value >= 1000 else f"${value}K"
    add_text(s, bar_area_left + bar_w + Inches(0.1), row_top,
             Inches(1.5), row_h,
             val_str, size=11, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)

# === RIGHT: By Cost Driver ===
right_x = Inches(6.83)
right_top = Inches(2.15)
right_w = Inches(6.0)
right_h = Inches(4.7)
add_rect(s, right_x, right_top, right_w, right_h, SAND)
add_text(s, right_x + Inches(0.25), right_top + Inches(0.15),
         right_w - Inches(0.5), Inches(0.35),
         "BY COST DRIVER (USD)", size=12, bold=True, color=NAVY)

drivers = [
    ("Consultants & Pro Services", 1910, CLAY),
    ("Staff Salaries & Benefits", 1710, NAVY),
    ("Indirect Cost Allocations", 472, SKY),
    ("Programmatic Delivery", 62, MUTED),
    ("Travel & Accommodation", 57, MUTED),
    ("Ops & Admin Support", 24, MUTED),
]
max_drv = max(v for _, v, _ in drivers)

drv_left = right_x + Inches(2.1)
drv_top = right_top + Inches(0.65)
drv_w = right_w - Inches(2.7)
drv_h = Inches(0.50)
drv_gap = Inches(0.13)

for i, (label, value, color) in enumerate(drivers):
    row_top = drv_top + i * (drv_h + drv_gap)
    add_text(s, right_x + Inches(0.20), row_top,
             Inches(1.85), drv_h,
             label, size=10, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
    bar_w = int(drv_w * (value / max_drv))
    add_rect(s, drv_left, row_top, bar_w, drv_h, color, corner=False)
    val_str = f"${value/1000:.2f}M" if value >= 1000 else f"${value}K"
    add_text(s, drv_left + bar_w + Inches(0.08), row_top,
             Inches(1.2), drv_h,
             val_str, size=10, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)

# Bottom takeaway
ta_top = Inches(6.95)
add_text(s, Inches(0.5), ta_top, Inches(12.3), Inches(0.3),
         "≈85% of spend went to consultants and staff. "
         "The product has yet to reach commercial release.",
         size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

add_brand_bar(s)


# ============================================================================
# SLIDE 4 — Market Reality: Why Selling ERP to NGOs Is Hard
# ============================================================================

s = new_slide()
slide_title_bar(s, "Market Reality — Why Selling ERP to NGOs Is Hard",
                "A saturated market of nonprofit accounting tools, sold at deep discount, "
                "with high switching costs.")

# Three columns
col_top = Inches(1.15)
col_h = Inches(5.55)
col_w = Inches(4.0)
col_gap = Inches(0.3)
col_xs = [Inches(0.5),
          Inches(0.5) + col_w + col_gap,
          Inches(0.5) + 2 * (col_w + col_gap)]

columns = [
    ("THE CROWDED MARKET",
     CLAY,
     "NGOs already have many strong, NGO-tuned options:",
     [
         ("QuickBooks Online", "Class tracking + fund accounting; free/discounted via TechSoup"),
         ("Aplos", "Built for nonprofits; ~USD 60–190/mo"),
         ("Sage Intacct for NPOs", "High-end, dimension-based fund accounting"),
         ("Blackbaud Financial Edge NXT", "Long-standing nonprofit ERP standard"),
         ("MIP Fund Accounting", "Community Brands; strong in mid-size INGOs"),
         ("NetSuite for Nonprofits", "Oracle SuiteSuccess SKU"),
         ("Xero · Zoho Books · Wave · Odoo Community", "Low-cost / free SMB tier"),
     ]),
    ("WHY NGOs WON'T SWITCH",
     NAVY,
     "Even when a better option exists, the math rarely works:",
     [
         ("Already on something", "Migration cost ≫ marginal feature gain"),
         ("Indirect-cost caps", "Donors cap overheads at 7–15%; software competes with salaries"),
         ("Donor reporting", "Funders demand familiar exports — Blackbaud, Sage, QB"),
         ("Country-by-country complexity", "Local tax, payroll, statutory reporting — hard to centralise"),
         ("Capacity to operate", "Many lack the finance team to run any ERP at all"),
     ]),
    ("THE PRICING SQUEEZE",
     SKY,
     "Incumbents have neutralised price as a wedge:",
     [
         ("TechSoup discounts", "60–90% off list for QB, Sage, Microsoft, Blackbaud"),
         ("Google for Nonprofits", "Free workspace + Workspace Marketplace add-ons"),
         ("Microsoft for Nonprofits", "Free / steeply discounted M365 + Dynamics 365 grants"),
         ("Free tier ERPs", "Odoo Community, ERPNext — same engine, USD 0 licence"),
         ("Bundled finance services", "Big-4 / regional consultancies bundle their own stacks"),
     ]),
]

for x, (title, accent, intro, items) in zip(col_xs, columns):
    add_rect(s, x, col_top, col_w, col_h, SAND)
    add_rect(s, x, col_top, col_w, Inches(0.55), accent, corner=False)
    add_text(s, x + Inches(0.2), col_top + Inches(0.05),
             col_w - Inches(0.4), Inches(0.45),
             title, size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), col_top + Inches(0.65),
             col_w - Inches(0.4), Inches(0.45),
             intro, size=10, italic=True, color=MUTED)
    add_bullets(s, x + Inches(0.25), col_top + Inches(1.15),
                col_w - Inches(0.5), col_h - Inches(1.3),
                items, size=10, bullet_color=accent)

# Bottom takeaway band
ta_top = Inches(6.80)
add_rect(s, Inches(0.5), ta_top, Inches(12.3), Inches(0.45), CLAY_SOFT)
add_text(s, Inches(0.7), ta_top + Inches(0.04),
         Inches(11.9), Inches(0.4),
         "Takeaway: even an NGO-customised ERP is a feature competing in a "
         "commoditised category against deeply discounted, donor-trusted incumbents.",
         size=11, bold=True, italic=True, color=CLAY,
         anchor=MSO_ANCHOR.MIDDLE)

add_brand_bar(s)


# ============================================================================
# SLIDE 5 — The Strategic Pivot
# ============================================================================

s = new_slide()
slide_title_bar(s, "The Strategic Pivot",
                "Keep the grant platform as the wedge. Re-channel the ERP to a buyer "
                "who can actually fund it.")

# THEN vs NOW header bars
then_x = Inches(0.5)
now_x = Inches(6.83)
band_w = Inches(6.0)

add_rect(s, then_x, Inches(1.1), band_w, Inches(0.5), MUTED, corner=False)
add_text(s, then_x, Inches(1.1), band_w, Inches(0.5),
         "THEN — Original Strategy",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, now_x, Inches(1.1), band_w, Inches(0.5), CLAY, corner=False)
add_text(s, now_x, Inches(1.1), band_w, Inches(0.5),
         "NOW — Refocused Strategy",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# THEN content
add_rect(s, then_x, Inches(1.65), band_w, Inches(2.4), SAND)
then_items = [
    ("Two products in parallel", "KujaLink + Kuja ERP, both sold direct"),
    ("Buyer = NGOs and donors", "Both products pushed to both sides"),
    ("ERP as recurring revenue", "Bet on NGOs paying licences"),
    ("Back-office service alongside", "Bundled where needed"),
]
add_bullets(s, then_x + Inches(0.25), Inches(1.80),
            band_w - Inches(0.5), Inches(2.2),
            then_items, size=11, bullet_color=MUTED)

# NOW content
add_rect(s, now_x, Inches(1.65), band_w, Inches(2.4), CLAY_SOFT)
now_items = [
    ("Lead with grant management", "Kuja Grant Platform = the wedge product"),
    ("Two-sided market", "NGOs (capacity + visibility) + donors (oversight + AI)"),
    ("ERP repositioned", "Sold to donors / INGOs as a service THEY fund for grantees"),
    ("Services + ERP bundled for donors", "Capacity-building line item in grants"),
]
add_bullets(s, now_x + Inches(0.25), Inches(1.80),
            band_w - Inches(0.5), Inches(2.2),
            now_items, size=11, bullet_color=CLAY)

# Two GTM motions
gtm_top = Inches(4.25)
gtm_h = Inches(2.45)
add_text(s, Inches(0.5), gtm_top - Inches(0.05),
         Inches(12.3), Inches(0.3),
         "TWO DISTINCT GO-TO-MARKET MOTIONS",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

motions = [
    (Inches(0.5),
     "Kuja Grant Platform (SaaS)",
     NAVY,
     [
         "Sold to NGOs + donors directly",
         "Subscription + per-grant pricing",
         "Network effects: more donors → more NGOs → more donors",
         "Already shipped and live in production",
     ]),
    (Inches(6.83),
     "Kuja ERP-as-a-Service (Donor-funded)",
     CLAY,
     [
         "Sold to donors / INGOs, not to NGO buyers",
         "Donor sponsors ERP + back-office service for selected grantees",
         "Funded from existing capacity-building grant line items",
         "Aligns Odoo 17 customisation with who values grant budgeting most",
     ]),
]

for x, title, accent, bullets in motions:
    add_rect(s, x, gtm_top + Inches(0.3), band_w, gtm_h, SAND)
    add_rect(s, x, gtm_top + Inches(0.3), band_w, Inches(0.5),
             accent, corner=False)
    add_text(s, x, gtm_top + Inches(0.3), band_w, Inches(0.5),
             title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.25), gtm_top + Inches(0.90),
                band_w - Inches(0.5), gtm_h - Inches(0.65),
                bullets, size=11, bullet_color=accent)

# Bottom rationale band
ra_top = Inches(6.95)
add_text(s, Inches(0.5), ra_top, Inches(12.3), Inches(0.3),
         "Why this works: donors have budget for grantee capacity-building. "
         "NGOs don't have budget for ERP licences. Match the product to the wallet.",
         size=11, italic=True, color=CLAY,
         align=PP_ALIGN.CENTER)

add_brand_bar(s)


# ============================================================================
# SLIDE 6 — The problem we address
# ============================================================================

s = new_slide()
slide_title_bar(s, "The Problem We Address",
                "Grant funding for Global South non-profits is structurally inefficient.")

# Three columns
col_w = Inches(4.0); col_top = Inches(1.3); col_h = Inches(5.6)
col_gap = Inches(0.3)
xs = [Inches(0.5), Inches(0.5) + col_w + col_gap,
      Inches(0.5) + 2 * (col_w + col_gap)]

cols = [
    ("FOR NGOs", [
        "Capacity assessment fatigue — answer the same 80–120 questions across CHS, UN-HACT, NUPAS, STEP, and donor-specific variants",
        "Application overhead with no help to know if their answer is strong",
        "Compliance evidence scattered across email, drives, and binders",
        "Reporting templates that vary by donor with no guidance on what each donor cares about",
        "Decision opacity — most applications rejected with little or no feedback",
        "Language exclusion — donor portals are English-only",
    ]),
    ("FOR DONORS", [
        "Sourcing is slow and intuition-driven — finding strong NGOs requires personal networks",
        "Capacity signal is non-standard — every donor maintains their own framework",
        "Application quality variance comes from language proficiency, not programme quality",
        "Compliance is reactive — sanctions, registry, document validation happen at the last minute",
        "Reports do not roll up — aggregating outcomes requires manual extraction",
        "No early-warning signal when grants drift off track",
    ]),
    ("FOR REVIEWERS", [
        "Each application takes hours to read carefully",
        "Rationales drift — reviewers anchor on each other, not the rubric",
        "Patterns across the portfolio are invisible",
        "No standardised way to cite specific evidence per criterion",
        "Suspect quality cannot be flagged with confidence",
    ]),
]

for (title, items), x in zip(cols, xs):
    # Header
    add_rect(s, x, col_top, col_w, Inches(0.55), NAVY)
    add_text(s, x, col_top, col_w, Inches(0.55),
             title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body card
    add_rect(s, x, col_top + Inches(0.6), col_w, col_h - Inches(0.6),
             SAND, line_color=NAVY)
    add_bullets(s, x + Inches(0.18), col_top + Inches(0.75),
                col_w - Inches(0.35), col_h - Inches(0.8),
                items, size=10.5, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 7 — Market opportunity (with chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Market Opportunity",
                "$2.75B today · $4.79B by 2030 · 10.3% CAGR")

# Chart on the left
add_chart(s, "market_size", Inches(0.4), Inches(1.2), Inches(8.4))

# Key takeaways on the right
add_rect(s, Inches(9.0), Inches(1.2), Inches(3.9), Inches(5.7),
         SAND, line_color=NAVY)
add_text(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.4),
         "Why it matters now",
         size=13, bold=True, color=NAVY)
add_bullets(s, Inches(9.2), Inches(1.9), Inches(3.6), Inches(4.9),
            [
                ("Sector disruption. ",
                 "USAID dissolution in 2025 with 83% of programmes cancelled has created urgent demand for diversified funding infrastructure."),
                ("AI inflection. ",
                 "92% of nonprofits use AI in some capacity, but 76% lack formal AI strategy — there is a window for an AI-native platform."),
                ("Cloud shift. ",
                 "65% of new grant management deployments are cloud-based."),
                ("Localisation pressure. ",
                 "Grand Bargain 2.0 targets 25% direct local funding; actual <5%."),
            ],
            size=10, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 8 — TAM / SAM / SOM
# ============================================================================

s = new_slide()
slide_title_bar(s, "Total / Serviceable / Obtainable Market",
                "Three views on the same opportunity")
add_chart(s, "tam_sam_som", Inches(0.7), Inches(1.2), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 9 — Sector trends (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "The Sector Is Shifting",
                "Four trends are creating urgent demand for new grant infrastructure")
add_chart(s, "sector_trends", Inches(0.7), Inches(1.2), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 10 — What Kuja Is
# ============================================================================

s = new_slide()
slide_title_bar(s, "What Kuja Is",
                "An end-to-end platform anchored on three category-defining commitments.")

# Three columns describing the commitments
col_w = Inches(4.0); col_top = Inches(1.3); col_h = Inches(5.6)
col_gap = Inches(0.3)
xs = [Inches(0.5), Inches(0.5) + col_w + col_gap,
      Inches(0.5) + 2 * (col_w + col_gap)]

commits = [
    ("TWO-PILLAR TRUST PROFILE",
     "Every NGO presents a Capacity Profile (\"can they execute?\") and a Due Diligence Profile (\"are they safe to fund?\"). Both pillars passport across applications.",
     ["Five capacity frameworks built in",
      "Custom donor framework builder",
      "Sanctions / AML / CTF",
      "Registration & standing",
      "Beneficial ownership + COI",
      "Adverse media monitoring"]),
    ("EMBEDDED AI INTELLIGENCE",
     "Twenty-one AI surfaces woven into every workflow. Action-oriented, grounded in user data, traceable to sources, editable end to end.",
     ["Application & report drafting",
      "Real-time document scoring",
      "Pre-submission readiness",
      "Reviewer one-screen summary",
      "Compliance pre-empt + pre-flight",
      "Portfolio insights"]),
    ("COMPLIANCE FOR BOTH SIDES",
     "Simplifies compliance and risk for the entire relationship — NGOs are supported through every obligation; donors get pre-assessed reports + portfolio health.",
     ["AI compliance calendar",
      "Proactive reminders",
      "AI report drafting",
      "4-pillar health score",
      "Trajectory + slips forecast",
      "Risk register with owners"]),
]

for (title, blurb, items), x in zip(commits, xs):
    # Header
    add_rect(s, x, col_top, col_w, Inches(0.55), CLAY)
    add_text(s, x, col_top, col_w, Inches(0.55),
             title, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body card
    add_rect(s, x, col_top + Inches(0.6), col_w, col_h - Inches(0.6),
             CLAY_SOFT, line_color=CLAY)
    add_text(s, x + Inches(0.2), col_top + Inches(0.75),
             col_w - Inches(0.4), Inches(1.5),
             blurb, size=10.5, color=NAVY, italic=True)
    add_bullets(s, x + Inches(0.2), col_top + Inches(2.2),
                col_w - Inches(0.4), col_h - Inches(2.4),
                items, size=10.5, color=NAVY)

add_brand_bar(s)


# ============================================================================
# SLIDE 11 — Two-Pillar Trust Profile (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "The Organisation Trust Profile",
                "Two pillars · one source of truth · travels with the organisation")
add_chart(s, "trust_profile", Inches(0.7), Inches(1.2), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 12 — Capacity Passporting
# ============================================================================

s = new_slide()
slide_title_bar(s, "Capacity Passporting",
                "Do the assessment once. Carry it forward to every donor's framework.")

# Visual: 1 → many
# Source on left
add_rect(s, Inches(0.7), Inches(2.8), Inches(3.2), Inches(2.2),
         CLAY_SOFT, line_color=CLAY)
add_text(s, Inches(0.7), Inches(3.0), Inches(3.2), Inches(0.5),
         "STEP 1", size=11, bold=True, color=CLAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(3.5), Inches(3.2), Inches(1.5),
         "NGO completes the Kuja Capacity Framework once · stored on the org profile",
         size=12.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Middle: passport
add_rect(s, Inches(4.5), Inches(2.5), Inches(3.5), Inches(2.8),
         NAVY)
add_text(s, Inches(4.5), Inches(2.7), Inches(3.5), Inches(0.5),
         "STEP 2", size=11, bold=True, color=CLAY_SOFT, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.5), Inches(3.2), Inches(3.5), Inches(0.7),
         "Capacity Passport",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.5), Inches(3.9), Inches(3.5), Inches(1.3),
         "Framework-neutral facts · AI maps across target frameworks · provenance traced to source",
         size=11, color=CLAY_SOFT, italic=True, align=PP_ALIGN.CENTER)

# Right: target frameworks
add_text(s, Inches(8.7), Inches(2.5), Inches(4.0), Inches(0.4),
         "STEP 3 — Prefilled across any donor framework",
         size=11, bold=True, color=CLAY)
target_frameworks = ["UN-HACT", "STEP", "CHS", "NUPAS", "Donor custom"]
for i, fw in enumerate(target_frameworks):
    y = Inches(2.95 + i * 0.55)
    add_rect(s, Inches(8.7), y, Inches(4.0), Inches(0.45),
             SAND, line_color=NAVY)
    add_text(s, Inches(8.7), y, Inches(4.0), Inches(0.45),
             fw, size=11, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Big effort-reduction stat at the bottom
add_rect(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(1.2),
         CLAY)
add_text(s, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.4),
         "Effort reduction per subsequent application",
         size=11, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6),
         "12 hours of duplicated work  →  under 1 hour of review and refinement",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_brand_bar(s)


# ============================================================================
# SLIDE 13 — Due Diligence — four sub-pillars
# ============================================================================

s = new_slide()
slide_title_bar(s, "Due Diligence Beyond Sanctions",
                "Four sub-pillars · passported across applications · auditable end to end")

dd = [
    ("Sanctions, AML & CTF",
     "UN · OFAC · EU + watchlists + CTF designations.",
     "Continuous re-screening on schedule, on change, and on every new application. Coverage includes declared beneficial owners and named officers."),
    ("Registration & Standing",
     "7 government registries today · expanding to all jurisdictions with public portals.",
     "Legal name, status, number verified against organisation-claimed values. Tax-exempt certification verified per country."),
    ("Beneficial Ownership",
     "Officer disclosure + ownership-chain transparency.",
     "Conflict-of-interest checks against donor staff directories surface conflicts before award. Historical preservation for retrospective audit."),
    ("Adverse Media",
     "Daily news scan with AI-classified relevance and severity.",
     "Auto-raised risk register entries above configured thresholds. Plain-language context summary with source link and suggested next steps."),
]
# 2x2 grid
card_w = Inches(6.1); card_h = Inches(2.7)
positions = [
    (Inches(0.5), Inches(1.2)),
    (Inches(6.75), Inches(1.2)),
    (Inches(0.5), Inches(4.05)),
    (Inches(6.75), Inches(4.05)),
]
for (title, lead, body), (x, y) in zip(dd, positions):
    add_rect(s, x, y, card_w, card_h, SAND, line_color=CLAY)
    # Title bar
    add_rect(s, x, y, card_w, Inches(0.5), CLAY)
    add_text(s, x + Inches(0.2), y, card_w - Inches(0.3), Inches(0.5),
             title, size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Lead line
    add_text(s, x + Inches(0.2), y + Inches(0.6),
             card_w - Inches(0.3), Inches(0.6),
             lead, size=11.5, bold=True, color=CLAY, italic=True)
    # Body
    add_text(s, x + Inches(0.2), y + Inches(1.2),
             card_w - Inches(0.3), card_h - Inches(1.3),
             body, size=11, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 14 — Embedded AI Intelligence
# ============================================================================

s = new_slide()
slide_title_bar(s, "Embedded AI Intelligence",
                "21 AI-driven surfaces woven into every workflow — grounded, traceable, action-oriented.")

# Three columns for the three primary workspaces
col_w = Inches(4.0); col_top = Inches(1.3); col_h = Inches(5.7)
col_gap = Inches(0.3)
xs = [Inches(0.5), Inches(0.5) + col_w + col_gap,
      Inches(0.5) + 2 * (col_w + col_gap)]

cols = [
    ("NGO", [
        "Match scoring",
        "Application co-author",
        "Real-time document scoring",
        "Submission readiness",
        "Compliance pre-empt",
        "Report co-author",
        "Donor-perspective pre-flight",
        "Compliance to-do list",
    ]),
    ("DONOR", [
        "Grant brief from 2-line prompt",
        "Grant import from PDF/DOCX/TXT",
        "Median-NGO preview",
        "Burden critique",
        "Portfolio insights",
        "Cross-grant patterns",
        "Compliance explanation",
        "Health narrative",
    ]),
    ("REVIEWER", [
        "One-screen summary",
        "Evidence extraction",
        "Comparable signal",
        "Decision-changers",
        "Per-criterion rationale",
        "Red-flag detection",
        "Suggest-criteria",
        "Audit-ready record",
    ]),
]

for (title, items), x in zip(cols, xs):
    add_rect(s, x, col_top, col_w, Inches(0.55), NAVY)
    add_text(s, x, col_top, col_w, Inches(0.55),
             title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 8 chips per column, 0.55in each + 0.08in gap
    for i, item in enumerate(items):
        y_chip = col_top + Inches(0.7) + Inches(i * 0.6)
        add_rect(s, x + Inches(0.1), y_chip, col_w - Inches(0.2),
                 Inches(0.5), CLAY_SOFT, line_color=CLAY)
        add_text(s, x + Inches(0.2), y_chip, col_w - Inches(0.4),
                 Inches(0.5), item, size=11, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

add_brand_bar(s)


# ============================================================================
# SLIDE 15 — Compliance simplified for both sides (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Compliance — Simplified for Both Sides",
                "The platform's flagship: NGOs are supported, donors are confident.")
add_chart(s, "compliance_both_sides", Inches(0.7), Inches(1.2), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 16 — Mobile + Offline
# ============================================================================

s = new_slide()
slide_title_bar(s, "Mobile-First, Offline-First",
                "The operating reality of Global South CSOs — not an afterthought.")

# Big stat block on the left
add_rect(s, Inches(0.5), Inches(1.4), Inches(5.5), Inches(5.5),
         NAVY)
add_text(s, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.6),
         "WHY IT MATTERS", size=12, bold=True, color=CLAY_SOFT,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(2.3), Inches(5.1), Inches(2.5),
         "The primary device for most Global South NGO staff is a phone. "
         "Connectivity is intermittent. Bandwidth is limited. The platform "
         "treats this as a first-class scenario rather than a degraded mode.",
         size=13, color=WHITE, italic=True)
add_text(s, Inches(0.7), Inches(5.0), Inches(5.1), Inches(1.8),
         "Field officers draft, capture evidence, and complete reports "
         "without a connection — the system syncs in the background when the "
         "network returns.",
         size=13, color=WHITE)

# Capabilities grid on the right
caps = [
    ("Progressive Web App",
     "Installable on phone home screen · same login as desktop"),
    ("Offline-first",
     "Draft applications, complete assessments, write reports without a network"),
    ("Continuous auto-save",
     "Every keystroke saved to local device storage"),
    ("Background sync",
     "Queued uploads transmit automatically on reconnection"),
    ("Camera as evidence",
     "Photographs of consent forms, beneficiaries, before/after — vision-based extraction"),
    ("Low-bandwidth optimised",
     "Heavy AI work runs async; core flows work on slow 3G"),
]
for i, (k, v) in enumerate(caps):
    row = i // 2; col = i % 2
    x = Inches(6.3 + col * 3.4); y = Inches(1.4 + row * 1.9)
    add_rect(s, x, y, Inches(3.2), Inches(1.7), CLAY_SOFT, line_color=CLAY)
    add_text(s, x + Inches(0.15), y + Inches(0.15),
             Inches(2.9), Inches(0.4),
             k, size=11.5, bold=True, color=CLAY)
    add_text(s, x + Inches(0.15), y + Inches(0.55),
             Inches(2.9), Inches(1.1),
             v, size=10, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 17 — Competitive landscape (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Competitive Landscape",
                "No competitor combines Global South focus with deep AI integration.")
add_chart(s, "competitive_2x2", Inches(0.7), Inches(1.1), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 18 — Competitors are mostly donor-centric
# ============================================================================

s = new_slide()
slide_title_bar(s, "Most Competitors Are Donor-Centric",
                "The NGO is the structurally underserved side of the market.")

# Three category cards across the top (60% of slide)
cat_top = Inches(1.2); cat_h = Inches(4.0); cat_gap = Inches(0.2)
card_w = (SLIDE_W - Inches(1.0) - 2 * cat_gap) / 3

categories = [
    ("DONOR-CENTRIC GRANT MANAGEMENT",
     "Optimised for the funder's workflow. NGOs see a submission portal — not a working partner.",
     NAVY,
     [
         ("Fluxx",        "Enterprise; deep donor lifecycle"),
         ("Submittable",  "Proven scale (25M apps)"),
         ("SmartSimple",  "Configurable, 45+ languages"),
         ("Foundant GLM", "Donor-friendly, user-friendly"),
         ("Benevity",     "Corporate CSR, $14B+ donations"),
         ("Bonterra",     "Grantmaker product launched Feb 2026"),
     ]),
    ("NGO-DISCOVERY ONLY",
     "Tools that help NGOs find grants — but stop there. No application, review, reporting, or compliance.",
     SKY,
     [
         ("Instrumentl",  "450K+ funders, 31K+ RFPs"),
         ("OpenGrants",   "AI matching + writer marketplace"),
         ("GrantHub",     "Affordable pipeline tracker"),
     ]),
    ("TWO-SIDED / SPECIALIST",
     "Rare. None combine marketplace, trust profile, embedded AI, and compliance support.",
     CLAY,
     [
         ("GlobalGiving",       "Marketplace, 175 countries — no GM"),
         ("UN Partner Portal",  "UN ↔ CSO only, no AI"),
         ("TechSoup STEP",      "Assessment only, no GM"),
         ("Xapien",             "Due diligence specialist, no GM"),
     ]),
]

for i, (cat, blurb, color, comps) in enumerate(categories):
    x = Inches(0.5) + i * (card_w + cat_gap)
    # Header
    add_rect(s, x, cat_top, card_w, Inches(0.55), color)
    add_text(s, x, cat_top, card_w, Inches(0.55), cat,
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Body card
    add_rect(s, x, cat_top + Inches(0.6), card_w, cat_h - Inches(0.6),
             SAND, line_color=color)
    add_text(s, x + Inches(0.2), cat_top + Inches(0.75),
             card_w - Inches(0.35), Inches(0.85),
             blurb, size=10, color=NAVY, italic=True)
    # Competitor list
    for j, (name, desc) in enumerate(comps):
        row_y = cat_top + Inches(1.7) + Inches(j * 0.4)
        add_text(s, x + Inches(0.2), row_y, Inches(1.5), Inches(0.35),
                 "• " + name, size=10, bold=True, color=color)
        add_text(s, x + Inches(1.7), row_y,
                 card_w - Inches(1.85), Inches(0.35),
                 desc, size=9, color=INK, italic=True)

# Big bottom callout — Kuja's positioning
callout_top = Inches(5.4); callout_h = Inches(1.65)
add_rect(s, Inches(0.5), callout_top, SLIDE_W - Inches(1.0), callout_h,
         CLAY)
add_text(s, Inches(0.7), callout_top + Inches(0.18),
         SLIDE_W - Inches(1.4), Inches(0.5),
         "KUJA — THE ONLY DUAL-SIDED PLATFORM WITH NGO-FIRST AI",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), callout_top + Inches(0.7),
         SLIDE_W - Inches(1.4), Inches(0.95),
         "Marketplace · Two-pillar Organisation Trust Profile · Embedded AI Intelligence working partner "
         "for NGOs AND donors · Compliance simplified for both sides · Mobile + offline · ERP. "
         "No competitor combines these — most never tried to serve the NGO at all.",
         size=11, color=WHITE, italic=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_brand_bar(s)


# ============================================================================
# SLIDE 19 — Core Feature Matrix (table)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Core Feature Matrix",
                "Y = fully supported · P = partial · — = not available. Kuja column highlighted.")

# Build the table natively in pptx.
# Rows: 1 header + core feature rows
# Cols: 1 feature name + 8 competitor cols
competitors_cols = ["Kuja", "Fluxx", "Submit-\ntable", "Smart-\nSimple",
                    "Bonterra", "UN Partner\nPortal", "Global-\nGiving", "Instrumentl"]
# Pick the 18 most decisive features (the differentiators that tell the story)
feature_rows = [
    # (Feature, [Kuja, Fluxx, Submittable, SmartSimple, Bonterra, UNPP, GlobalGiving, Instrumentl])
    ("Embedded AI Intelligence (action-oriented)", ["Y", "P", "P", "P", "P", "—", "—", "P"]),
    ("Two-Pillar Organisation Trust Profile",      ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Capacity Passporting (one assessment, many frameworks)",
                                                   ["Y", "—", "—", "—", "—", "P", "—", "—"]),
    ("Custom Donor Framework Builder",             ["Y", "P", "P", "Y", "P", "—", "—", "—"]),
    ("AI Application Co-Author with Provenance",   ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Real-Time Document Scoring on Upload",       ["Y", "P", "P", "P", "P", "—", "—", "—"]),
    ("Pre-Submission Readiness + Compliance Pre-Empt",
                                                   ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Reviewer One-Screen Summary + Evidence Extraction",
                                                   ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Sanctions / AML / CTF Screening",            ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Govt. Registry & Standing Verification",     ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Beneficial Ownership + Adverse Media",       ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("AI Compliance Calendar + Reminders",         ["Y", "—", "—", "P", "—", "—", "—", "—"]),
    ("AI Report Drafting + Donor Pre-Flight",      ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("4-Pillar Health + Slips Forecast",           ["Y", "—", "—", "—", "—", "—", "—", "—"]),
    ("Two-Sided Marketplace + Match Scoring",      ["Y", "—", "—", "—", "—", "P", "Y", "Y"]),
    ("Mobile-First PWA + Offline-First",           ["Y", "—", "—", "—", "—", "P", "Y", "Y"]),
    ("Six-Language UI with Role-Aware Tone",       ["Y", "P", "—", "Y", "—", "Y", "P", "—"]),
    ("Audit Chain + 2FA + GDPR-Erasure",           ["Y", "Y", "Y", "Y", "P", "Y", "—", "—"]),
]

n_rows = 1 + len(feature_rows)
n_cols = 1 + len(competitors_cols)

# Table dimensions
table_left = Inches(0.4); table_top = Inches(1.05)
table_width = SLIDE_W - Inches(0.8)
table_height = Inches(6.0)

table_shape = s.shapes.add_table(n_rows, n_cols,
                                 table_left, table_top,
                                 table_width, table_height)
table = table_shape.table

# Set column widths: feature column wider, competitor columns equal.
# python-pptx requires integer EMU values, so we int() the division results.
feature_col_w = Inches(4.2)
comp_col_w = int((table_width - feature_col_w) / len(competitors_cols))
table.columns[0].width = feature_col_w
for ci in range(1, n_cols):
    table.columns[ci].width = comp_col_w

# Set row heights: header taller for 2-line competitor names
header_h = Inches(0.55)
row_h = int((table_height - header_h) / len(feature_rows))
table.rows[0].height = header_h
for ri in range(1, n_rows):
    table.rows[ri].height = row_h


def _style_cell(cell, text, *, bold=False, font_size=8.5, fg=INK,
                bg=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = fg


# Header row
_style_cell(table.cell(0, 0), "Feature", bold=True, font_size=10,
            fg=WHITE, bg=NAVY, align=PP_ALIGN.LEFT)
for ci, comp in enumerate(competitors_cols, start=1):
    bg = CLAY if comp == "Kuja" else NAVY
    _style_cell(table.cell(0, ci), comp, bold=True, font_size=9.5,
                fg=WHITE, bg=bg, align=PP_ALIGN.CENTER)

# Data rows
for ri, (feat, vals) in enumerate(feature_rows, start=1):
    row_bg = SAND if ri % 2 == 1 else WHITE
    _style_cell(table.cell(ri, 0), feat, bold=False, font_size=8.5,
                fg=INK, bg=row_bg, align=PP_ALIGN.LEFT)
    for ci, v in enumerate(vals, start=1):
        if v == "Y":
            fg = WHITE if competitors_cols[ci - 1] == "Kuja" else GROW
            bold = True
        elif v == "P":
            fg = WHITE if competitors_cols[ci - 1] == "Kuja" else SUN
            bold = True
        else:
            fg = WHITE if competitors_cols[ci - 1] == "Kuja" else MUTED
            bold = False
        # Kuja column gets clay shading; other Y cells get a light green hint
        if competitors_cols[ci - 1] == "Kuja":
            bg = CLAY
        else:
            bg = row_bg
        _style_cell(table.cell(ri, ci), v, bold=bold,
                    font_size=11 if v in ("Y", "P") else 10,
                    fg=fg, bg=bg, align=PP_ALIGN.CENTER)

# Below-table takeaway
add_text(s, Inches(0.4), Inches(7.05), SLIDE_W - Inches(0.8), Inches(0.22),
         "Kuja delivers every feature fully. Closest competitor (Fluxx / SmartSimple) "
         "covers fewer than half — and none on the differentiators specific to Global South NGOs.",
         size=9.5, color=CLAY, italic=True, align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 20 — Feature comparison (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Feature Coverage vs. Competitors",
                "40/40 fully supported · 17 differentiators no competitor offers.")
add_chart(s, "feature_count", Inches(0.7), Inches(1.2), Inches(11.9))
add_brand_bar(s)


# ============================================================================
# SLIDE 21 — Competitive Advantages (category-defining)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Competitive Advantages — Category-Defining",
                "Three commitments no competitor matches.")

# Three big numbered cards
cards = [
    ("1", "Embedded AI Intelligence, Not a Chatbot",
     "AI is woven through every workflow as a working partner. Grounded in user data, traceable to sources, action-oriented (an editable starting point), and falls back gracefully when the AI service is unavailable. Competitors offering AI today provide chatbots or analytics dashboards — not a working AI partner."),
    ("2", "Two-Pillar Organisation Trust Profile",
     "Every NGO presents Capacity + Due Diligence side by side, both passported across applications. No competitor offers both pillars together. No competitor offers passporting on the due-diligence side."),
    ("3", "Compliance Simplified for Both Sides",
     "NGOs get an AI compliance working partner (extract obligations → calendar → reminders → drafting → scoring → pre-flight → revision guidance). Donors get pre-assessed scored reports, 4-pillar health, trajectory forecasting, slips-in-N-days warnings, and a risk register. The donor is never surprised; the NGO is never blocked."),
]
card_h = Inches(1.9); card_top = Inches(1.2); gap = Inches(0.15)
for i, (num, title, body) in enumerate(cards):
    y = card_top + i * (card_h + gap)
    add_rect(s, Inches(0.5), y, Inches(12.3), card_h, SAND, line_color=CLAY)
    # Number badge
    add_rect(s, Inches(0.5), y, Inches(1.2), card_h, CLAY)
    add_text(s, Inches(0.5), y, Inches(1.2), card_h,
             num, size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.9), y + Inches(0.2), Inches(10.7), Inches(0.5),
             title, size=15, bold=True, color=CLAY)
    add_text(s, Inches(1.9), y + Inches(0.7), Inches(10.7), Inches(1.1),
             body, size=11, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 22 — Competitive Advantages (structural)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Competitive Advantages — Structural",
                "Seven additional structural strengths.")

# 4 cards on left, 3 on right (or 3+4) — 7 cards in a 4×2 layout
adv = [
    ("4", "Only End-to-End Platform",
     "Marketplace, trust profile, AI applications, dual-scoring review, AI-supported reporting + health, ERP — single platform."),
    ("5", "Mobile-First, Offline-First",
     "Progressive web app · auto-save · background sync · vision-based evidence extraction."),
    ("6", "Built By and For the Global South",
     "Six languages at parity with role-aware tone · Adeso's 30+ years of humanitarian experience."),
    ("7", "Dual-Sided Network Effects",
     "CSOs gain visibility, AI working partner, calendar of obligations. Donors gain vetted pipeline + portfolio health."),
    ("8", "Auditable by Construction",
     "Provenance on every AI claim · tamper-evident hash-chained audit log · donor override transparency."),
    ("9", "Seamless ERP Conversion",
     "KujaBuild (Odoo 17) onboards seamlessly because most data is already captured in earlier stages."),
    ("10", "Credibility of Adeso",
     "Somali-founded · NEAR Network co-founder · three decades of African humanitarian impact."),
]
# Lay out 4 wide x 2 high (one empty)
card_w = Inches(3.0); card_h = Inches(2.7); gap_x = Inches(0.2); gap_y = Inches(0.2)
start_x = Inches(0.5); start_y = Inches(1.2)
for i, (num, title, body) in enumerate(adv):
    row = i // 4; col = i % 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(s, x, y, card_w, card_h, SAND, line_color=NAVY)
    # Number circle
    add_rect(s, x + Inches(0.15), y + Inches(0.15),
             Inches(0.6), Inches(0.6), CLAY)
    add_text(s, x + Inches(0.15), y + Inches(0.15),
             Inches(0.6), Inches(0.6),
             num, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.9), y + Inches(0.2),
             card_w - Inches(1.05), Inches(0.7),
             title, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.15), y + Inches(1.0),
             card_w - Inches(0.3), card_h - Inches(1.1),
             body, size=10.5, color=INK)

add_brand_bar(s)


# ============================================================================
# SLIDE 23 — Target market segments
# ============================================================================

s = new_slide()
slide_title_bar(s, "Target Market Segments",
                "Five segments · phased geographic expansion · dual-sided model.")

segs = [
    ("African CSOs / NGOs", "100,000+ registered",
     "Free marketplace tier + paid assessment + ERP option", NAVY),
    ("Global South CSOs (LatAm, Asia)", "500,000+",
     "Phase 2–3 geographic expansion", SKY),
    ("Bilateral / multilateral donors", "50+ major institutions",
     "Premium end-to-end solution", CLAY),
    ("Private foundations", "200,000+ globally",
     "Application portal + matching", CLAY),
    ("INGOs seeking local partners", "1,000+",
     "Marketplace + assessment integration", SKY),
]
card_w = Inches(11.5); card_top = Inches(1.3); card_h = Inches(0.95)
gap = Inches(0.1)
for i, (seg, size, approach, color) in enumerate(segs):
    y = card_top + i * (card_h + gap)
    # Big card
    add_rect(s, Inches(0.9), y, card_w, card_h, SAND, line_color=color)
    # Coloured stripe
    add_rect(s, Inches(0.9), y, Inches(0.18), card_h, color, corner=False)
    # Segment name
    add_text(s, Inches(1.3), y + Inches(0.05),
             Inches(4.5), card_h - Inches(0.1),
             seg, size=13, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    # Size
    add_text(s, Inches(5.8), y + Inches(0.05),
             Inches(2.5), card_h - Inches(0.1),
             size, size=12.5, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE)
    # Approach
    add_text(s, Inches(8.3), y + Inches(0.05),
             Inches(4.0), card_h - Inches(0.1),
             approach, size=10.5, color=INK, italic=True,
             anchor=MSO_ANCHOR.MIDDLE)

add_brand_bar(s)


# ============================================================================
# SLIDE 24 — Roadmap (chart)
# ============================================================================

s = new_slide()
slide_title_bar(s, "Roadmap",
                "v5.0 launch wave shipped Q2 2026 · momentum continues through 2028.")
add_chart(s, "roadmap", Inches(0.5), Inches(1.1), Inches(12.3))
add_brand_bar(s)


# ============================================================================
# SLIDE 25 — Why Adeso, Why Now (closing)
# ============================================================================

s = new_slide()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY, corner=False)
add_rect(s, Inches(0), Inches(3.6), SLIDE_W, Inches(0.06), CLAY, corner=False)

# Title
add_text(s, Inches(0), Inches(0.7), SLIDE_W, Inches(0.7),
         "Why Adeso · Why Now",
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(1.4), SLIDE_W, Inches(0.5),
         "Right team · right moment · right product.",
         size=14, color=CLAY_SOFT, align=PP_ALIGN.CENTER, italic=True)

# Three pillars
pillars = [
    ("WHY ADESO",
     "Somali-founded African social enterprise. 30+ years of humanitarian impact. "
     "Co-founder of the NEAR Network — the largest network of national and local "
     "organisations in the Global South. Built by the communities Kuja serves."),
    ("WHY NOW",
     "USAID dissolution in 2025 ($36B cut, 83% of programmes cancelled) has created "
     "urgent demand for diversified funding infrastructure and direct donor-to-NGO "
     "connections. 92% of nonprofits already use AI but 76% lack a strategy."),
    ("WHY KUJA",
     "The only platform combining marketplace, two-pillar Trust Profile, Embedded AI "
     "Intelligence, compliance simplified for both sides, mobile-first delivery, and "
     "ERP operations — purpose-built for the Global South."),
]
for i, (title, body) in enumerate(pillars):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(4.1), Inches(4.0), Inches(2.7),
             WHITE)
    # Title bar
    add_rect(s, x, Inches(4.1), Inches(4.0), Inches(0.55), CLAY)
    add_text(s, x, Inches(4.1), Inches(4.0), Inches(0.55),
             title, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), Inches(4.8),
             Inches(3.6), Inches(2.0),
             body, size=11, color=INK)

# Closing footer
add_text(s, Inches(0), Inches(7.1), SLIDE_W, Inches(0.3),
         "Findable. Fundable. Trusted.",
         size=14, italic=True, color=CLAY_SOFT, align=PP_ALIGN.CENTER)


# ============================================================================
# SAVE
# ============================================================================

out_path = os.path.join(THIS_DIR, "Kuja_Market_Analysis_Deck_v5.pptx")
prs.save(out_path)

import os
size_kb = os.path.getsize(out_path) / 1024
size_str = f"{size_kb:.1f} KB" if size_kb < 1024 else f"{size_kb / 1024:.1f} MB"
print(f"Generated: {out_path}")
print(f"File size: {size_str}")
print(f"Slides: {len(prs.slides)}")
