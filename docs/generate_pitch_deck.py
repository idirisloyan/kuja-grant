"""
Kuja Partner Pitch Deck Generator
Generates a professional 16-slide PowerPoint presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
# CONSTANTS & BRANDING
# ──────────────────────────────────────────────
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)       # Deep blue
SECONDARY = RGBColor(0x0D, 0x94, 0x88)     # Teal
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)        # Amber
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x1F, 0x2A, 0x37)      # Near-black for body
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)    # Subtle bg
MED_GRAY = RGBColor(0x9C, 0xA3, 0xAF)      # Muted text
GREEN_CHECK = RGBColor(0x16, 0xA3, 0x4A)    # Green for check
RED_X = RGBColor(0xDC, 0x26, 0x26)          # Red for X
PARTIAL_YLW = RGBColor(0xF5, 0x9E, 0x0B)   # Amber for partial

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ──────────────────────────────────────────────
# HELPER FUNCTIONS
# ──────────────────────────────────────────────
def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_NAME, word_wrap=True, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.15:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rich_textbox(slide, left, top, width, height):
    """Return textframe for multi-paragraph content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_TEXT,
                  alignment=PP_ALIGN.LEFT, space_after=6, space_before=0,
                  bullet=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, font_color=WHITE, bold=False,
                     alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    if text:
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = FONT_NAME
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=12, font_color=WHITE, bold=False,
             alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    if text:
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = FONT_NAME
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


def add_arrow_right(slide, left, top, width, height, fill_color=SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, fill_color=SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text="",
               font_size=11, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = FONT_NAME
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
    return shape


def add_slide_number(slide, num, total=16):
    add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text, slide_num, total=16):
    """Add a consistent blue header bar with title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0), SLIDE_W, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.7),
                 title_text, font_size=32, bold=True, color=WHITE)
    # Teal accent line below header
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0), Inches(1.0), SLIDE_W, Inches(0.05))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = SECONDARY
    accent_line.line.fill.background()
    add_slide_number(slide, slide_num, total)


def add_bottom_bar(slide):
    """Add subtle bottom brand bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(7.15), SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_text_box(slide, Inches(0.5), Inches(7.15), Inches(4), Inches(0.35),
                 "KUJA  |  Powered by Adeso", font_size=9, color=WHITE, bold=True)


def add_callout_box(slide, left, top, width, height, text, font_size=13,
                    fill_color=ACCENT, font_color=DARK_TEXT):
    shape = add_rounded_rect(slide, left, top, width, height, fill_color,
                             text, font_size, font_color, bold=True)
    return shape


# ──────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ──────────────────────────────────────────────
def create_slide_1():
    slide = prs.slides.add_slide(blank_layout)
    # Full blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Decorative teal accent strip on left
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = SECONDARY
    strip.line.fill.background()

    # Decorative geometric elements - right side circles (subtle)
    for i, (x, y, s, alpha) in enumerate([
        (10.5, 0.5, 2.0, 0.08),
        (11.0, 5.0, 1.5, 0.06),
        (9.0, 6.0, 1.0, 0.05),
    ]):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x), Inches(y), Inches(s), Inches(s))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x3B, 0x5E, 0xC4)  # Slightly lighter blue
        c.line.fill.background()

    # KUJA title
    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(10), Inches(1.5),
                 "KUJA", font_size=80, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)

    # Amber accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1.0), Inches(3.1), Inches(2.5), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Subtitle
    add_text_box(slide, Inches(1.0), Inches(3.4), Inches(10), Inches(1.0),
                 "The First AI-Powered End-to-End\nGrant Management Platform for the Global South",
                 font_size=28, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Tagline
    add_text_box(slide, Inches(1.0), Inches(4.8), Inches(10), Inches(0.5),
                 "Powered by Adeso \u2014 30 Years of African-Led Development",
                 font_size=18, bold=False, color=RGBColor(0xA5, 0xB4, 0xFC),
                 alignment=PP_ALIGN.LEFT)

    # Date
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(4), Inches(0.4),
                 "March 2026", font_size=16, color=RGBColor(0x93, 0xA3, 0xBF),
                 alignment=PP_ALIGN.LEFT)

    # Website
    add_text_box(slide, Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.4),
                 "kuja.org  |  adesoafrica.org", font_size=14,
                 color=RGBColor(0x93, 0xA3, 0xBF), alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 1)


# ──────────────────────────────────────────────
# SLIDE 2: THE PROBLEM
# ──────────────────────────────────────────────
def create_slide_2():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "The Funding Gap", 2)
    add_bottom_bar(slide)

    # Three columns
    col_w = Inches(3.6)
    col_h = Inches(3.8)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    top_y = Inches(1.4)

    labels = ["FOR CSOs", "FOR DONORS", "SYSTEMIC"]
    colors = [PRIMARY, SECONDARY, ACCENT]
    contents = [
        [
            "\u2022  4\u201312 duplicate assessments per year",
            "\u2022  Below 5% direct funding to local orgs",
            "\u2022  English-first platforms exclude millions",
            "\u2022  Enterprise pricing excludes local orgs",
        ],
        [
            "\u2022  Cannot find vetted local partners",
            "\u2022  Risk perception barriers persist",
            "\u2022  Manual compliance monitoring",
            "\u2022  Lack of standardized assessments",
        ],
        [
            "\u2022  $36B USAID cuts (2025)",
            "\u2022  Aid localization stalled at <5%",
            "\u2022  No infrastructure for direct partnerships",
            "\u2022  Growing demand, shrinking supply",
        ],
    ]

    for i in range(3):
        x = start_x + i * (col_w + gap)
        # Card background
        card = add_rounded_rect(slide, x, top_y, col_w, col_h, LIGHT_GRAY)
        # Label bar at top of card
        label_bar = add_rounded_rect(slide, x + Inches(0.15), top_y + Inches(0.15),
                                     col_w - Inches(0.3), Inches(0.55), colors[i],
                                     labels[i], font_size=16, font_color=WHITE, bold=True)
        # Content
        tf = add_rich_textbox(slide, x + Inches(0.25), top_y + Inches(0.9),
                              col_w - Inches(0.5), col_h - Inches(1.0))
        for j, line in enumerate(contents[i]):
            add_paragraph(tf, line, font_size=14, color=DARK_TEXT,
                          space_after=8, space_before=4 if j > 0 else 0)

    # Bottom callout
    add_callout_box(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.7),
                    "Local organizations closest to communities have the most effective solutions \u2014 but the least resources.",
                    font_size=15, fill_color=ACCENT, font_color=DARK_TEXT)


# ──────────────────────────────────────────────
# SLIDE 3: OUR SOLUTION
# ──────────────────────────────────────────────
def create_slide_3():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "One Platform. Complete Grant Lifecycle.", 3)
    add_bottom_bar(slide)

    stages = [
        ("1", "DISCOVER", "Marketplace"),
        ("2", "ASSESS", "Capacity"),
        ("3", "VERIFY", "Due Diligence"),
        ("4", "MATCH", "AI Matching"),
        ("5", "APPLY", "AI-Guided"),
        ("6", "REVIEW", "Dual Score"),
        ("7", "REPORT", "AI Compliance"),
        ("8", "MANAGE", "ERP"),
    ]

    # Flow diagram - 8 stages in two rows
    box_w = Inches(1.35)
    box_h = Inches(1.1)
    arrow_w = Inches(0.35)
    arrow_h = Inches(0.3)
    start_x = Inches(0.5)
    row1_y = Inches(1.5)
    row2_y = Inches(3.5)

    for idx, (num, name, desc) in enumerate(stages[:4]):
        x = start_x + idx * (box_w + arrow_w + Inches(0.15))
        # Number circle
        add_circle(slide, x + box_w / 2 - Inches(0.2), row1_y - Inches(0.25),
                   Inches(0.4), SECONDARY, num, font_size=14)
        # Box
        shape = add_rounded_rect(slide, x, row1_y + Inches(0.2), box_w, box_h,
                                 PRIMARY, "", font_size=14, font_color=WHITE, bold=True)
        tf = shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = FONT_NAME
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER

        # Arrow (except after last in row)
        if idx < 3:
            add_arrow_right(slide, x + box_w + Inches(0.05),
                            row1_y + Inches(0.55), arrow_w, arrow_h, SECONDARY)

    # Down arrow from row 1 to row 2 (right side)
    down_arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(6.8), row1_y + box_h + Inches(0.25),
                                        Inches(0.3), Inches(0.5))
    down_arrow.fill.solid()
    down_arrow.fill.fore_color.rgb = SECONDARY
    down_arrow.line.fill.background()

    # Second row (stages 5-8, right to left visually but we'll do left to right with adjusted order)
    # Actually let's keep it flowing: down then continue left to right in row 2
    for idx, (num, name, desc) in enumerate(stages[4:]):
        x = start_x + idx * (box_w + arrow_w + Inches(0.15))
        # Number circle
        add_circle(slide, x + box_w / 2 - Inches(0.2), row2_y - Inches(0.25),
                   Inches(0.4), SECONDARY, num, font_size=14)
        # Box
        shape = add_rounded_rect(slide, x, row2_y + Inches(0.2), box_w, box_h,
                                 PRIMARY, "", font_size=14, font_color=WHITE, bold=True)
        tf = shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = FONT_NAME
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
        p2.font.name = FONT_NAME
        p2.alignment = PP_ALIGN.CENTER

        if idx < 3:
            add_arrow_right(slide, x + box_w + Inches(0.05),
                            row2_y + Inches(0.55), arrow_w, arrow_h, SECONDARY)

    # Closing arrow loops back (decorative curved line from stage 8 back)
    # We'll add a subtle return arrow
    ret_arrow = slide.shapes.add_shape(MSO_SHAPE.CURVED_LEFT_ARROW,
                                       start_x + 3 * (box_w + arrow_w + Inches(0.15)) + box_w + Inches(0.15),
                                       row2_y + Inches(0.3),
                                       Inches(0.6), Inches(0.7))
    ret_arrow.fill.solid()
    ret_arrow.fill.fore_color.rgb = ACCENT
    ret_arrow.line.fill.background()

    # Non-ERP reporting note between rows
    add_rounded_rect(slide, Inches(8.3), row2_y + Inches(0.15), Inches(4.5), Inches(1.2),
                     LIGHT_GRAY)
    tf_note = add_rich_textbox(slide, Inches(8.5), row2_y + Inches(0.2), Inches(4.1), Inches(1.1))
    add_paragraph(tf_note, "WITHOUT OUR ERP?", font_size=12, bold=True, color=SECONDARY, space_after=2)
    add_paragraph(tf_note, "Stages 1\u20137 work standalone. NGOs not using the ERP use our built-in reporting module to seamlessly submit compliance reports per grant agreement requirements.", font_size=10, color=DARK_TEXT, space_after=0)

    # Bottom statement
    add_callout_box(slide, Inches(3.0), Inches(5.5), Inches(7.3), Inches(0.6),
                    "No other platform covers all 8 stages.", font_size=18,
                    fill_color=ACCENT, font_color=DARK_TEXT)


# ──────────────────────────────────────────────
# SLIDE 4: ADESO'S UNIQUE ADVANTAGE
# ──────────────────────────────────────────────
def create_slide_4():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Built by Those Who Know the Challenge", 4)
    add_bottom_bar(slide)

    # Left column - Adeso Credentials
    left_x = Inches(0.6)
    col_w = Inches(5.5)

    add_rounded_rect(slide, left_x, Inches(1.3), col_w, Inches(0.55), SECONDARY,
                     "ADESO CREDENTIALS", font_size=16, font_color=WHITE, bold=True)

    credentials = [
        "\u2022  Founded 1991 in Somalia \u2014 30+ years humanitarian experience",
        "\u2022  Led the NEAR network for aid localization advocacy",
        "\u2022  Operations across Kenya, Somalia, and the Horn of Africa",
        "\u2022  501(c)(3) through Myriad USA",
        "\u2022  Deep relationships: Oxfam, IRC, Save the Children,",
        "   Gates Foundation, Hilton Foundation, Porticus",
    ]
    tf = add_rich_textbox(slide, left_x + Inches(0.2), Inches(2.05), col_w - Inches(0.4), Inches(3.5))
    for line in credentials:
        add_paragraph(tf, line, font_size=15, color=DARK_TEXT, space_after=8)

    # Right column - Why This Matters
    right_x = Inches(6.8)

    add_rounded_rect(slide, right_x, Inches(1.3), col_w, Inches(0.55), PRIMARY,
                     "WHY THIS MATTERS", font_size=16, font_color=WHITE, bold=True)

    matters = [
        "\u2022  Understands compliance from BOTH sides",
        "   \u2014 as an NGO AND as a grantor",
        "\u2022  Built by Global South, for Global South",
        "\u2022  Multi-language: Arabic, French, Spanish",
        "   live today \u2014 no competitor offers this",
        "\u2022  Not a Silicon Valley startup guessing",
        "   \u2014 lived experience solving these problems",
    ]
    tf2 = add_rich_textbox(slide, right_x + Inches(0.2), Inches(2.05), col_w - Inches(0.4), Inches(3.5))
    for line in matters:
        add_paragraph(tf2, line, font_size=15, color=DARK_TEXT, space_after=8)

    # Bottom insight bar
    add_callout_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.65),
                    "\"We didn't build Kuja because we read about the problem. We built it because we lived it.\"",
                    font_size=15, fill_color=LIGHT_GRAY, font_color=PRIMARY)


# ──────────────────────────────────────────────
# SLIDE 5: CAPACITY ASSESSMENT
# ──────────────────────────────────────────────
def create_slide_5():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "AI-Powered Capacity Assessment", 5)
    add_bottom_bar(slide)

    # Center hub
    center_x = Inches(4.0)
    center_y = Inches(2.8)
    hub_w = Inches(2.4)
    hub_h = Inches(1.2)
    add_rounded_rect(slide, center_x, center_y, hub_w, hub_h, PRIMARY,
                     "Kuja\nAssessment\nEngine", font_size=16, font_color=WHITE, bold=True)

    # 5 current framework nodes + extensible node arranged around the hub
    frameworks = [
        ("Kuja\nFramework", Inches(1.2), Inches(1.5), SECONDARY),
        ("STEP\n(TechSoup)", Inches(6.0), Inches(1.5), SECONDARY),
        ("UN-HACT", Inches(7.2), Inches(3.5), SECONDARY),
        ("CHS\nAlliance", Inches(5.5), Inches(5.2), SECONDARY),
        ("NUPAS", Inches(1.5), Inches(4.5), SECONDARY),
        ("+ Donor\nSpecific", Inches(0.5), Inches(3.0), ACCENT),
    ]

    fw_w = Inches(1.6)
    fw_h = Inches(1.0)
    for name, fx, fy, fw_color in frameworks:
        add_rounded_rect(slide, fx, fy, fw_w, fw_h, fw_color,
                         name, font_size=13, font_color=WHITE, bold=True)
        # Connector line from framework to hub center
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      min(fx + fw_w / 2, center_x + hub_w / 2),
                                      min(fy + fw_h / 2, center_y + hub_h / 2),
                                      Inches(0.04),
                                      abs((fy + fw_h / 2) - (center_y + hub_h / 2)) or Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        line.line.fill.background()

    # Right side - Key benefits
    benefits_x = Inches(9.0)
    add_rounded_rect(slide, benefits_x, Inches(1.3), Inches(3.8), Inches(0.5),
                     PRIMARY, "KEY BENEFITS", font_size=14, font_color=WHITE, bold=True)

    benefits = [
        ("Self-service", "Complete online, no expensive consultants"),
        ("Passportable", "Assess once, share with multiple donors"),
        ("AI-analyzed", "Upload policies, AI identifies gaps and scores"),
        ("Extensible", "Add donor-specific or any new framework on demand"),
        ("Configurable", "Donors set their own criteria and weightings"),
    ]

    for i, (title, desc) in enumerate(benefits):
        y = Inches(2.0) + i * Inches(1.0)
        add_circle(slide, benefits_x, y, Inches(0.4), ACCENT, str(i + 1), font_size=12)
        tf = add_rich_textbox(slide, benefits_x + Inches(0.55), y - Inches(0.05),
                              Inches(3.2), Inches(0.85))
        add_paragraph(tf, title, font_size=14, bold=True, color=PRIMARY)
        add_paragraph(tf, desc, font_size=12, color=DARK_TEXT, space_after=0)


# ──────────────────────────────────────────────
# SLIDE 6: AI CAPABILITIES
# ──────────────────────────────────────────────
def create_slide_6():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "AI That Works Across the Entire Lifecycle", 6)
    add_bottom_bar(slide)

    capabilities = [
        ("Document\nAnalysis", "Score 0\u2013100 with\ndetailed findings"),
        ("Grant Agreement\nParsing", "Extract requirements\nautomatically"),
        ("Report\nEvaluation", "Per-requirement\ncompliance scoring"),
        ("Capacity\nAssessment", "AI reviews policies,\nidentifies gaps"),
        ("Grant\nMatching", "Smart NGO-donor\nmatching algorithm"),
        ("Chat\nAssistant", "Role-aware guidance\n(NGO vs Donor)"),
    ]

    box_w = Inches(3.5)
    box_h = Inches(1.6)
    start_x = Inches(0.8)
    gap_x = Inches(0.5)
    row1_y = Inches(1.5)
    row2_y = Inches(3.6)

    icons = ["\U0001F4C4", "\U0001F4CB", "\u2705", "\U0001F50D", "\U0001F517", "\U0001F4AC"]

    for i, (title, desc) in enumerate(capabilities):
        row = i // 3
        col = i % 3
        x = start_x + col * (box_w + gap_x)
        y = row1_y if row == 0 else row2_y

        # Card
        card = add_rounded_rect(slide, x, y, box_w, box_h, LIGHT_GRAY)
        # Color accent at top of card
        accent_bar = add_rect(slide, x, y, box_w, Inches(0.06), PRIMARY if col % 2 == 0 else SECONDARY)

        # Number badge
        add_circle(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.45),
                   PRIMARY if col % 2 == 0 else SECONDARY,
                   str(i + 1), font_size=14, font_color=WHITE)

        # Title
        add_text_box(slide, x + Inches(0.7), y + Inches(0.15), box_w - Inches(0.9), Inches(0.6),
                     title.replace("\n", " "), font_size=16, bold=True, color=PRIMARY)

        # Description
        add_text_box(slide, x + Inches(0.7), y + Inches(0.7), box_w - Inches(0.9), Inches(0.8),
                     desc.replace("\n", " "), font_size=13, color=DARK_TEXT)

    # Bottom bar
    add_rounded_rect(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.55),
                     PRIMARY, "Powered by Anthropic Claude  |  5 languages live (EN, FR, AR, ES, SW)  \u2014  expanding to 10+",
                     font_size=15, font_color=WHITE, bold=True)


# ──────────────────────────────────────────────
# SLIDE 7: LIVE DUE DILIGENCE
# ──────────────────────────────────────────────
def create_slide_7():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Real-Time Compliance Verification", 7)
    add_bottom_bar(slide)

    # Left section - Government Registry
    left_x = Inches(0.5)
    sec_w = Inches(5.8)

    add_rounded_rect(slide, left_x, Inches(1.4), sec_w, Inches(0.55), SECONDARY,
                     "GOVERNMENT REGISTRY VERIFICATION",
                     font_size=15, font_color=WHITE, bold=True)

    # Current 7 countries
    countries = [
        ("Kenya", "NGO Board / BRS"),
        ("Nigeria", "Corporate Affairs Commission"),
        ("South Africa", "DSD NPO Directorate"),
        ("Uganda", "NGO Bureau"),
        ("Tanzania", "NiS Registry"),
        ("Somalia", "MOIFAR"),
        ("Ethiopia", "ACSO"),
    ]

    for i, (country, registry) in enumerate(countries):
        row = i % 4
        col = i // 4
        x = left_x + Inches(0.2) + col * Inches(2.8)
        y = Inches(2.2) + row * Inches(0.55)
        # Country pill
        add_rounded_rect(slide, x, y, Inches(2.5), Inches(0.42), LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.03), Inches(2.3), Inches(0.36),
                     f"{country} \u2014 {registry}", font_size=11, color=DARK_TEXT, bold=False)

    # Expandable message
    add_text_box(slide, left_x + Inches(0.2), Inches(4.55), sec_w - Inches(0.4), Inches(0.7),
                 "7 countries live today. Designed to expand to any Global South country\nwith a government verification portal.",
                 font_size=12, bold=True, color=SECONDARY, alignment=PP_ALIGN.LEFT)

    # Right section - Sanctions Screening
    right_x = Inches(6.8)

    add_rounded_rect(slide, right_x, Inches(1.4), sec_w, Inches(0.55), PRIMARY,
                     "SANCTIONS SCREENING  \u2014  4+ Databases",
                     font_size=15, font_color=WHITE, bold=True)

    sanctions = [
        ("UN Security Council", "Global sanctions list"),
        ("US OFAC SDN", "US Treasury sanctions"),
        ("EU Financial Sanctions", "European Union list"),
        ("World Bank Debarment", "Procurement exclusions"),
        ("OpenSanctions API", "Unified global database"),
    ]

    for i, (name, desc) in enumerate(sanctions):
        y = Inches(2.2) + i * Inches(0.65)
        add_rounded_rect(slide, right_x + Inches(0.2), y, Inches(5.2), Inches(0.5), LIGHT_GRAY)
        tf = add_rich_textbox(slide, right_x + Inches(0.35), y + Inches(0.02),
                              Inches(4.8), Inches(0.45))
        p = add_paragraph(tf, "", font_size=12)
        run1 = p.add_run()
        run1.text = name
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = PRIMARY
        run1.font.name = FONT_NAME
        run2 = p.add_run()
        run2.text = f"  \u2014  {desc}"
        run2.font.size = Pt(12)
        run2.font.color.rgb = DARK_TEXT
        run2.font.name = FONT_NAME

    # Bottom callout
    add_callout_box(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(0.6),
                    "Live, automated, continuous \u2014 not a one-time check",
                    font_size=18, fill_color=ACCENT, font_color=DARK_TEXT)


# ──────────────────────────────────────────────
# SLIDE 8: FOR NGOs — THE JOURNEY
# ──────────────────────────────────────────────
def create_slide_8():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Invisible to Investment-Ready", 8)
    add_bottom_bar(slide)

    # Subtitle
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(0.4),
                 "THE NGO JOURNEY", font_size=14, bold=True, color=SECONDARY)

    steps = [
        ("1", "JOIN", "Free profile on\nkuja.org marketplace"),
        ("2", "ASSESS", "AI capacity assessment,\nget readiness score"),
        ("3", "IMPROVE", "30/60/90-day roadmap,\nlearning resources"),
        ("4", "APPLY", "AI-guided applications,\ndocument coaching"),
        ("5", "REPORT", "Seamless compliance\nreporting"),
    ]

    step_w = Inches(2.0)
    step_h = Inches(2.5)
    chevron_w = Inches(0.4)
    start_x = Inches(0.6)
    top_y = Inches(2.0)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (step_w + chevron_w + Inches(0.1))

        # Step card
        card = add_rounded_rect(slide, x, top_y, step_w, step_h, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1.5)

        # Number circle at top
        add_circle(slide, x + step_w / 2 - Inches(0.3), top_y + Inches(0.15),
                   Inches(0.6), PRIMARY, num, font_size=20, font_color=WHITE)

        # Title
        add_text_box(slide, x + Inches(0.1), top_y + Inches(0.9),
                     step_w - Inches(0.2), Inches(0.4),
                     title, font_size=18, bold=True, color=PRIMARY,
                     alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x + Inches(0.1), top_y + Inches(1.4),
                     step_w - Inches(0.2), Inches(1.0),
                     desc.replace("\n", " "), font_size=13, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)

        # Chevron arrow between steps
        if i < 4:
            add_chevron(slide, x + step_w + Inches(0.02),
                        top_y + step_h / 2 - Inches(0.2),
                        chevron_w, Inches(0.4), SECONDARY)

    # Progress bar at bottom
    bar_y = Inches(5.0)
    bar_w = Inches(11.5)
    add_rect(slide, start_x, bar_y, bar_w, Inches(0.08), LIGHT_GRAY)
    # Gradient fill sections
    colors_prog = [
        RGBColor(0xDC, 0xED, 0xFC),
        RGBColor(0xBA, 0xDB, 0xF9),
        RGBColor(0x7C, 0xBE, 0xF3),
        RGBColor(0x3B, 0x82, 0xF6),
        RGBColor(0x1E, 0x40, 0xAF),
    ]
    seg_w = bar_w // 5
    for i in range(5):
        add_rect(slide, start_x + Emu(seg_w * i), bar_y, Emu(seg_w), Inches(0.08),
                 colors_prog[i])

    add_text_box(slide, Inches(0.6), Inches(5.3), Inches(11), Inches(0.4),
                 "Invisible  \u2192  Visible  \u2192  Assessed  \u2192  Ready  \u2192  Funded  \u2192  Compliant",
                 font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 9: FOR DONORS — THE VALUE
# ──────────────────────────────────────────────
def create_slide_9():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Searching to Certainty", 9)
    add_bottom_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5), Inches(0.4),
                 "THE DONOR JOURNEY", font_size=14, bold=True, color=SECONDARY)

    steps = [
        ("1", "DISCOVER", "Browse verified,\nassessed organizations"),
        ("2", "CONFIGURE", "Set your criteria,\nthresholds, weightings"),
        ("3", "REVIEW", "AI-ranked shortlists,\ndual scoring"),
        ("4", "AWARD", "Grant wizard with AI\nagreement parsing"),
        ("5", "MONITOR", "Real-time reporting,\nAI compliance scores"),
    ]

    step_w = Inches(2.0)
    step_h = Inches(2.5)
    chevron_w = Inches(0.4)
    start_x = Inches(0.6)
    top_y = Inches(2.0)

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (step_w + chevron_w + Inches(0.1))

        card = add_rounded_rect(slide, x, top_y, step_w, step_h, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1.5)

        add_circle(slide, x + step_w / 2 - Inches(0.3), top_y + Inches(0.15),
                   Inches(0.6), SECONDARY, num, font_size=20, font_color=WHITE)

        add_text_box(slide, x + Inches(0.1), top_y + Inches(0.9),
                     step_w - Inches(0.2), Inches(0.4),
                     title, font_size=18, bold=True, color=SECONDARY,
                     alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.1), top_y + Inches(1.4),
                     step_w - Inches(0.2), Inches(1.0),
                     desc.replace("\n", " "), font_size=13, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)

        if i < 4:
            add_chevron(slide, x + step_w + Inches(0.02),
                        top_y + step_h / 2 - Inches(0.2),
                        chevron_w, Inches(0.4), PRIMARY)

    # Bottom callout
    add_callout_box(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.65),
                    "For donors using our ERP: multi-tenant real-time financial visibility",
                    font_size=15, fill_color=LIGHT_GRAY, font_color=PRIMARY)


# ──────────────────────────────────────────────
# SLIDE 10: MARKET OPPORTUNITY
# ──────────────────────────────────────────────
def create_slide_10():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "$4.79B Market. Zero Solutions for the Global South.", 10)
    add_bottom_bar(slide)

    stats = [
        ("$2.75B", "Current Market (2024)", PRIMARY, "Grand View Research"),
        ("$4.79B", "Projected (2030)", SECONDARY, "10.3% CAGR"),
        ("58.2%", "NGO Segment Share", ACCENT, "Largest segment"),
        ("90%", "Nonprofits Wanting AI", PRIMARY, "CEP 2025 Report"),
        ("100K+", "African CSOs", SECONDARY, "Registered organizations"),
        ("<5%", "Direct Funding to Local Actors", RED_X, "vs 25% Grand Bargain target"),
    ]

    # 3 columns x 2 rows of stat cards with accent bars
    card_w = Inches(3.5)
    card_h = Inches(2.2)
    gap = Inches(0.5)
    start_x = Inches(0.9)

    for i, (value, label, color, source) in enumerate(stats):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap)
        y = Inches(1.5) + row * (card_h + Inches(0.35))

        # Card background
        card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1)

        # Color accent bar at top
        add_rect(slide, x, y, card_w, Inches(0.06), color)

        # Value (large)
        add_text_box(slide, x, y + Inches(0.2), card_w, Inches(1.0),
                     value, font_size=48, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, x + Inches(0.2), y + Inches(1.25), card_w - Inches(0.4), Inches(0.5),
                     label, font_size=15, bold=True, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)

        # Source
        add_text_box(slide, x + Inches(0.2), y + Inches(1.7), card_w - Inches(0.4), Inches(0.4),
                     source, font_size=10, bold=False, color=MED_GRAY,
                     alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 11: COMPETITIVE LANDSCAPE
# ──────────────────────────────────────────────
def create_slide_11():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "No One Else Does What We Do", 11)
    add_bottom_bar(slide)

    # Comparison matrix - expanded to match product analysis doc
    features = [
        "End-to-End Lifecycle",
        "AI-Powered",
        "Capacity Assessment",
        "Multi-Language",
        "Global South Design",
        "Marketplace",
        "Affordable for CSOs",
        "Due Diligence",
        "Reporting Module",
    ]

    competitors = ["Kuja", "Fluxx", "Submittable", "SmartSimple", "UNPP", "GlobalGiving", "Instrumentl", "Xapien"]

    # Y = full, P = partial, N = no
    # Columns:     Kuja  Fluxx  Subm   Smart  UNPP  GGiv  Instr  Xapien
    matrix = [
        ["Y",  "P",   "P",   "P",   "N",  "N",  "N",   "N"],   # End-to-End
        ["Y",  "P",   "P",   "P",   "N",  "N",  "P",   "Y"],   # AI-Powered
        ["Y",  "N",   "N",   "N",   "P",  "N",  "N",   "N"],   # Capacity Assessment
        ["Y",  "N",   "N",   "N",   "N",  "N",  "N",   "N"],   # Multi-Language (Global South langs)
        ["Y",  "N",   "N",   "P",   "Y",  "Y",  "N",   "N"],   # Global South
        ["Y",  "N",   "N",   "N",   "N",  "Y",  "N",   "N"],   # Marketplace
        ["Y",  "N",   "N",   "N",   "Y",  "Y",  "N",   "N"],   # Affordable
        ["Y",  "N",   "N",   "N",   "N",  "N",  "N",   "P"],   # Due Diligence
        ["Y",  "Y",   "P",   "Y",   "P",  "P",  "N",   "N"],   # Reporting
    ]

    table_x = Inches(0.3)
    table_y = Inches(1.35)
    feat_col_w = Inches(2.2)
    data_col_w = Inches(1.3)
    row_h = Inches(0.55)
    header_h = Inches(0.55)

    # Header row
    add_rect(slide, table_x, table_y, feat_col_w, header_h, PRIMARY,
             "Feature", font_size=13, font_color=WHITE, bold=True)

    for j, comp in enumerate(competitors):
        x = table_x + feat_col_w + j * data_col_w
        color = SECONDARY if j == 0 else PRIMARY
        add_rect(slide, x, table_y, data_col_w, header_h, color,
                 comp, font_size=11, font_color=WHITE, bold=True)

    # Data rows
    for i, feat in enumerate(features):
        y = table_y + header_h + i * row_h
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY

        # Feature name
        feat_shape = add_rect(slide, table_x, y, feat_col_w, row_h, bg,
                              feat, font_size=12, font_color=DARK_TEXT, bold=True,
                              alignment=PP_ALIGN.LEFT)
        feat_shape.text_frame.margin_left = Pt(10)

        for j, val in enumerate(matrix[i]):
            x = table_x + feat_col_w + j * data_col_w
            if val == "Y":
                symbol = "\u2713"
                sym_color = GREEN_CHECK
            elif val == "P":
                symbol = "\u25D0"
                sym_color = PARTIAL_YLW
            else:
                symbol = "\u2717"
                sym_color = RED_X

            cell = add_rect(slide, x, y, data_col_w, row_h, bg,
                            symbol, font_size=20, font_color=sym_color, bold=True)

    # Kuja column highlight border
    kuja_x = table_x + feat_col_w
    total_h = header_h + len(features) * row_h
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       kuja_x, table_y, data_col_w, total_h)
    highlight.fill.background()
    highlight.line.color.rgb = SECONDARY
    highlight.line.width = Pt(3)


# ──────────────────────────────────────────────
# SLIDE 12: BUSINESS MODEL
# ──────────────────────────────────────────────
def create_slide_12():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Sustainable Revenue, Subsidized Access", 12)
    add_bottom_bar(slide)

    tiers = [
        ("FREE", "Marketplace profile\nBasic capacity assessment\nGrant discovery feed",
         LIGHT_GRAY, DARK_TEXT, "$0"),
        ("APPLICATION\nPORTAL", "AI matching + assessment\nDonor-configured criteria\nShortlist generation",
         SECONDARY, WHITE, "$10K\u2013$20K/yr\nper donor"),
        ("END-TO-END", "Full fund management\nERP integration\nCustom AI workflows",
         PRIMARY, WHITE, "Custom\nPricing"),
    ]

    tier_w = Inches(3.5)
    tier_h = Inches(3.8)
    gap = Inches(0.5)
    start_x = Inches(1.1)
    top_y = Inches(1.4)

    for i, (name, features, bg_color, txt_color, price) in enumerate(tiers):
        x = start_x + i * (tier_w + gap)

        # Tier card
        card = add_rounded_rect(slide, x, top_y, tier_w, tier_h, bg_color if i > 0 else WHITE)
        if i == 0:
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            card.line.width = Pt(1.5)

        # Tier label
        label_color = MED_GRAY if i == 0 else WHITE
        add_text_box(slide, x, top_y + Inches(0.15), tier_w, Inches(0.4),
                     f"TIER {i + 1}", font_size=12, bold=True,
                     color=label_color, alignment=PP_ALIGN.CENTER)

        # Tier name
        name_color = DARK_TEXT if i == 0 else WHITE
        add_text_box(slide, x, top_y + Inches(0.5), tier_w, Inches(0.7),
                     name, font_size=20, bold=True,
                     color=name_color, alignment=PP_ALIGN.CENTER)

        # Price
        price_color = PRIMARY if i == 0 else ACCENT
        add_text_box(slide, x, top_y + Inches(1.3), tier_w, Inches(0.7),
                     price, font_size=18, bold=True,
                     color=price_color, alignment=PP_ALIGN.CENTER)

        # Features list
        feat_color = DARK_TEXT if i == 0 else RGBColor(0xE0, 0xF2, 0xFE)
        tf = add_rich_textbox(slide, x + Inches(0.3), top_y + Inches(2.1),
                              tier_w - Inches(0.6), Inches(1.5))
        for line in features.split("\n"):
            add_paragraph(tf, f"\u2713  {line}", font_size=13, color=feat_color, space_after=6)

    # Key principle
    add_callout_box(slide, Inches(0.8), Inches(5.5), Inches(7.5), Inches(0.65),
                    "Free marketplace + basic assessment for CSOs  |  Paid tiers for donors drive sustainability",
                    font_size=14, fill_color=ACCENT, font_color=DARK_TEXT)

    # Revenue target
    add_rounded_rect(slide, Inches(8.8), Inches(5.5), Inches(3.7), Inches(0.65),
                     PRIMARY, "2026 Target: $1M revenue + $1M fundraised",
                     font_size=14, font_color=WHITE, bold=True)


# ──────────────────────────────────────────────
# SLIDE 13: GO-TO-MARKET ROADMAP
# ──────────────────────────────────────────────
def create_slide_13():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "From Kenya to the World", 13)
    add_bottom_bar(slide)

    # Timeline
    phases = [
        ("Q1\u2013Q2 2026", "KENYA", "8 pipeline partners\n3 CSO Network MoUs",
         PRIMARY),
        ("Q2\u2013Q4 2026", "GLOBAL NETWORKS", "NEAR, WACSI, EACSOF,\nSANGONet + Oxfam/IRC/Save",
         SECONDARY),
        ("Q3 2026\u2013\nQ2 2027", "SUB-SAHARAN\nAFRICA", "24 countries\nContinental expansion",
         RGBColor(0x7C, 0x3A, 0xED)),  # Purple
        ("Q3\u2013Q4 2027", "LATIN AMERICA\n& CARIBBEAN", "15 countries\nRegional partnerships",
         ACCENT),
    ]

    # Horizontal timeline line
    line_y = Inches(3.0)
    line_start = Inches(0.8)
    line_end = Inches(12.5)
    timeline_bar = add_rect(slide, line_start, line_y, line_end - line_start, Inches(0.06), PRIMARY)

    phase_w = Inches(2.6)
    gap = Inches(0.35)

    for i, (time, name, desc, color) in enumerate(phases):
        x = line_start + Inches(0.2) + i * (phase_w + gap)

        # Timeline node
        add_circle(slide, x + phase_w / 2 - Inches(0.18), line_y - Inches(0.15),
                   Inches(0.35), color, "", font_size=10)

        # Time label above
        add_text_box(slide, x, Inches(1.6), phase_w, Inches(0.4),
                     time, font_size=13, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)

        # Phase card below
        card = add_rounded_rect(slide, x, Inches(3.4), phase_w, Inches(2.0), color,
                                "", font_size=14, font_color=WHITE, bold=True)
        tf = card.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = FONT_NAME
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(12)

        for line in desc.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)

    # Partners bar
    add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
                     LIGHT_GRAY,
                     "Partners: NEAR  |  WACSI  |  EACSOF  |  SANGONet  |  TechSoup  |  EPIC Africa  |  Pledge for Change",
                     font_size=14, font_color=DARK_TEXT, bold=True)


# ──────────────────────────────────────────────
# SLIDE 14: TRACTION & MILESTONES
# ──────────────────────────────────────────────
def create_slide_14():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Where We Are Today", 14)
    add_bottom_bar(slide)

    metrics = [
        ("600+", "Active Users on kuja.org"),
        ("200+", "Local Organizations Onboarded"),
        ("5", "Capacity Assessment Frameworks"),
        ("7", "Country Registry Verifications"),
        ("4+", "Sanctions Databases"),
        ("5+", "Languages Live \u2014 Expanding"),
    ]

    card_w = Inches(3.5)
    card_h = Inches(1.7)
    gap = Inches(0.5)
    start_x = Inches(0.9)

    for i, (value, label) in enumerate(metrics):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap)
        y = Inches(1.3) + row * (card_h + Inches(0.25))

        card = add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_GRAY)

        # Big number
        add_text_box(slide, x, y + Inches(0.1), card_w, Inches(0.8),
                     value, font_size=42, bold=True, color=PRIMARY,
                     alignment=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.4), Inches(0.6),
                     label, font_size=14, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)

    # Additional milestones as horizontal pills below the cards
    pill_y = Inches(5.0)
    pill_h = Inches(0.5)
    pill_gap = Inches(0.3)

    extras = [
        ("\u2713  Live AI platform in production", Inches(3.2)),
        ("\u2713  Multi-language: AR, FR, ES live \u2014 no competitor offers this", Inches(5.2)),
        ("\u2713  ERP (Odoo 17) in development", Inches(3.0)),
    ]
    pill_x = Inches(0.7)
    for text, pw in extras:
        add_rounded_rect(slide, pill_x, pill_y, pw, pill_h, PRIMARY,
                         text, font_size=13, font_color=WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)
        pill_x += pw + pill_gap


# ──────────────────────────────────────────────
# SLIDE 15: THE ASK
# ──────────────────────────────────────────────
def create_slide_15():
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Let's Transform Grant Management Together", 15)
    add_bottom_bar(slide)

    # What we're looking for
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.4),
                 "WHAT WE'RE LOOKING FOR", font_size=16, bold=True, color=PRIMARY)

    asks = [
        ("\u2794  Strategic Partners", "Organizations who share our localization vision and can co-create solutions"),
        ("\u2794  Pilot Donors", "Foundations willing to pilot the end-to-end grant management solution"),
        ("\u2794  CSO Networks", "Networks ready to onboard their member organizations at scale"),
        ("\u2794  Technology Partners", "Companies that can help us scale infrastructure and integrations"),
    ]

    for i, (title, desc) in enumerate(asks):
        y = Inches(2.0) + i * Inches(1.1)
        # Card
        card = add_rounded_rect(slide, Inches(0.8), y, Inches(6.5), Inches(0.9), LIGHT_GRAY)
        tf = card.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY
        p1.font.name = FONT_NAME
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.font.name = FONT_NAME
        p2.space_before = Pt(4)

    # Contact info - right side
    contact_x = Inches(8.0)
    contact_box = add_rounded_rect(slide, contact_x, Inches(1.4), Inches(4.5), Inches(4.5), PRIMARY)
    tf = contact_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(15)

    p = tf.paragraphs[0]
    p.text = "CONTACT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME
    p.space_after = Pt(20)

    lines = [
        ("Adeso", 22, True, WHITE),
        ("African Development Solutions", 14, False, RGBColor(0xA5, 0xB4, 0xFC)),
        ("", 12, False, WHITE),
        ("kuja.org", 18, True, ACCENT),
        ("adesoafrica.org", 16, False, RGBColor(0xA5, 0xB4, 0xFC)),
        ("", 12, False, WHITE),
        ("info@adesoafrica.org", 16, True, WHITE),
    ]

    for text, size, bold, color in lines:
        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(size)
        p2.font.bold = bold
        p2.font.color.rgb = color
        p2.font.name = FONT_NAME
        p2.space_after = Pt(6)


# ──────────────────────────────────────────────
# SLIDE 16: CLOSING
# ──────────────────────────────────────────────
def create_slide_16():
    slide = prs.slides.add_slide(blank_layout)

    # Full blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Teal accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = SECONDARY
    strip.line.fill.background()

    # Decorative circles
    for x, y, s in [(10.0, 0.3, 2.5), (11.5, 5.5, 1.5), (0.5, 5.8, 1.0)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(x), Inches(y), Inches(s), Inches(s))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x3B, 0x5E, 0xC4)
        c.line.fill.background()

    # KUJA
    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(10), Inches(1.5),
                 "KUJA", font_size=80, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Amber accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(5.2), Inches(3.1), Inches(3.0), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Tagline
    add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
                 "Matching Local Brilliance with Global Funding",
                 font_size=30, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0),
                 "The first AI-powered grant management platform\nbuilt for the Global South, by the Global South.",
                 font_size=20, bold=False, color=RGBColor(0xA5, 0xB4, 0xFC),
                 alignment=PP_ALIGN.CENTER)

    # Bottom
    add_text_box(slide, Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.5),
                 "kuja.org  |  adesoafrica.org  |  info@adesoafrica.org",
                 font_size=16, color=RGBColor(0x93, 0xA3, 0xBF),
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 16)


# ──────────────────────────────────────────────
# GENERATE ALL SLIDES
# ──────────────────────────────────────────────
print("Generating Kuja Partner Pitch Deck...")

create_slide_1()
print("  Slide  1/16: Title Slide")
create_slide_2()
print("  Slide  2/16: The Problem")
create_slide_3()
print("  Slide  3/16: Our Solution")
create_slide_4()
print("  Slide  4/16: Adeso's Advantage")
create_slide_5()
print("  Slide  5/16: Capacity Assessment")
create_slide_6()
print("  Slide  6/16: AI Capabilities")
create_slide_7()
print("  Slide  7/16: Live Due Diligence")
create_slide_8()
print("  Slide  8/16: NGO Journey")
create_slide_9()
print("  Slide  9/16: Donor Journey")
create_slide_10()
print("  Slide 10/16: Market Opportunity")
create_slide_11()
print("  Slide 11/16: Competitive Landscape")
create_slide_12()
print("  Slide 12/16: Business Model")
create_slide_13()
print("  Slide 13/16: Go-to-Market Roadmap")
create_slide_14()
print("  Slide 14/16: Traction & Milestones")
create_slide_15()
print("  Slide 15/16: The Ask")
create_slide_16()
print("  Slide 16/16: Closing")

# Save
output_path = r"C:\Users\IdirisLoyan\kuja-grant\docs\Kuja_Partner_Pitch_Deck_V2.pptx"
prs.save(output_path)
file_size = os.path.getsize(output_path)
print(f"\nDone! Saved to: {output_path}")
print(f"File size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
print(f"Total slides: {len(prs.slides)}")
